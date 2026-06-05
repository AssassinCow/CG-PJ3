"""Build the final 10-minute CG-Tutor presentation deck.

The deck is generated from current repository artifacts only:
current architecture diagrams are drawn with PowerPoint shapes and current
video thumbnails are extracted from outputs/<concept>/final.mp4. Historical
W2/W3 images are intentionally not used.
"""

from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SLIDES = ROOT / "slides"
ASSETS = SLIDES / "assets"
OUT = SLIDES / "CG-Tutor.pptx"

CONCEPTS = [
    ("affine_transformation", "仿射变换", "矩阵/坐标轴/物体运动"),
    ("forward_kinematics_chain", "正向运动学", "关节层级/末端执行器"),
    ("mirror_reflection", "镜面反射", "入射光/法线/反射角"),
    ("prism_dispersion_teaching", "棱镜色散", "内部光路/RGB 出射"),
    ("shape_morphing", "形状变形", "插值/连续动画"),
]

RESULT_ROWS = [
    ("affine", "0", "0.8125", "0/4", "7/12", "best_with_violations"),
    ("FK chain", "1", "0.7150", "7/1", "7/20", "best_with_violations"),
    ("mirror", "2", "0.7750", "0/2", "16/15", "best_with_violations"),
    ("prism", "1", "0.8388", "0/8", "1/18", "best_with_violations"),
    ("morphing", "0", "0.7625", "1/6", "3/14", "best_with_violations"),
]

NAVY = RGBColor(0x19, 0x2D, 0x3F)
BLUE = RGBColor(0x2F, 0x66, 0xB3)
TEAL = RGBColor(0x1E, 0x8A, 0x8A)
GREEN = RGBColor(0x2D, 0x6A, 0x4F)
ACCENT = RGBColor(0xB2, 0x3A, 0x3A)
ORANGE = RGBColor(0xCC, 0x7A, 0x29)
INK = RGBColor(0x24, 0x28, 0x2E)
GREY = RGBColor(0x6E, 0x75, 0x7E)
LIGHT = RGBColor(0xF3, 0xF5, 0xF7)
PALE_BLUE = RGBColor(0xE9, 0xF0, 0xFA)
PALE_GREEN = RGBColor(0xE8, 0xF3, 0xEC)
PALE_RED = RGBColor(0xFA, 0xEA, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def ensure_thumbnails() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    thumbs: dict[str, Path] = {}
    for concept, _, _ in CONCEPTS:
        video = ROOT / "outputs" / concept / "final.mp4"
        thumb = ASSETS / f"{concept}_thumb.png"
        thumbs[concept] = thumb
        if not video.exists():
            continue
        cmd = [
            "ffmpeg", "-y", "-ss", "00:00:06", "-i", str(video),
            "-frames:v", "1", "-vf", "scale=640:360", str(thumb),
        ]
        subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=False)
    return thumbs


def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, *, font_size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = frame.margin_right = Inches(0.05)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    for idx, line in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *, font_size=15, color=INK):
    return add_text(slide, x, y, w, h, "\n".join(f"•  {i}" for i in items),
                    font_size=font_size, color=color)


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.55, 0.30, 12.3, 0.58, title,
             font_size=27, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, 0.55, 0.88, 12.3, 0.35, subtitle,
                 font_size=13, color=GREY)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.28), Inches(12.25), Inches(0.035)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, page_num, total):
    add_text(slide, 0.55, 7.08, 8.0, 0.25,
             "CG-Tutor | Agentic Blender Teaching Animation", font_size=8, color=GREY)
    add_text(slide, 11.8, 7.08, 1.4, 0.25, f"{page_num} / {total}",
             font_size=8, color=GREY, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, *, fill=LIGHT, accent=NAVY):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = accent
    rect.line.width = Pt(0.8)
    add_text(slide, x + 0.14, y + 0.14, w - 0.28, 0.30, title,
             font_size=13, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.16, y + 0.50, w - 0.32, h - 0.60, body,
             font_size=10, color=INK, align=PP_ALIGN.CENTER)


def add_table(slide, x, y, headers, rows, widths, *, row_h=0.42, font_size=10):
    for r_idx, row in enumerate([headers, *rows]):
        x_pos = x
        for c_idx, cell in enumerate(row):
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos), Inches(y + r_idx * row_h),
                Inches(widths[c_idx]), Inches(row_h),
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = NAVY if r_idx == 0 else LIGHT
            rect.line.color.rgb = WHITE if r_idx == 0 else RGBColor(0xD0, 0xD6, 0xDD)
            rect.line.width = Pt(0.35)
            box = slide.shapes.add_textbox(
                Inches(x_pos), Inches(y + r_idx * row_h + 0.055),
                Inches(widths[c_idx]), Inches(row_h - 0.06),
            )
            para = box.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(cell)
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(font_size)
            run.font.bold = r_idx == 0
            run.font.color.rgb = WHITE if r_idx == 0 else INK
            x_pos += widths[c_idx]


def add_arrow(slide, x1, y1, x2, y2, *, color=GREY):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.4)
    line.line.end_arrowhead = True
    return line


def slide_title(prs, thumbs):
    s = add_blank(prs)
    add_text(s, 0.7, 0.85, 12.0, 0.75, "CG-Tutor",
             font_size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, 1.1, 1.72, 11.2, 0.55,
             "基于 Agent 的计算机图形学教学动画自动生成与验证框架",
             font_size=21, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, 1.0, 2.35, 11.4, 0.45,
             "Concept YAML → Storyboard → Blender scene.py → Render/Critic → final.mp4",
             font_size=15, color=GREY, align=PP_ALIGN.CENTER)
    x = 0.75
    for concept, zh, _ in CONCEPTS:
        thumb = thumbs.get(concept)
        if thumb and thumb.exists():
            s.shapes.add_picture(str(thumb), Inches(x), Inches(3.25), width=Inches(2.35))
        add_text(s, x, 5.65, 2.35, 0.45, zh, font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        x += 2.52
    add_text(s, 0.9, 6.45, 11.5, 0.35,
             "刘卓鑫 24300240170  |  朱劲舟 24300240117",
             font_size=13, color=GREEN, align=PP_ALIGN.CENTER)


def slide_requirements(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "选题与 PJ3 要求对应", "Agent + 图形学相关课题；报告/PPT/展示结果；强调创新性与实现")
    add_card(s, 0.7, 1.7, 3.8, 2.0, "我们的选题",
             "输入普通 CG 概念\n自动生成 3D 教学动画\n输出 MP4 + scene.py + 诊断 artifact",
             fill=PALE_BLUE, accent=BLUE)
    add_card(s, 4.85, 1.7, 3.8, 2.0, "Agent 属性",
             "Decomposer / Profile / Storyboard\nBlender Coder / Critic Ensemble\n多轮反馈和修复计划",
             fill=PALE_GREEN, accent=GREEN)
    add_card(s, 9.0, 1.7, 3.4, 2.0, "图形学属性",
             "坐标变换、运动学、反射、折射/色散、形状插值\n全部通过 Blender 渲染验证",
             fill=PALE_RED, accent=ACCENT)
    add_bullets(s, 0.9, 4.45, 11.8, 1.5, [
        "创新点不是“让 LLM 写一段 Blender 代码”，而是把成功状态变成可检查的中间产物。",
        "系统不把“看起来能渲染”直接当作成功；final_status 会保留 best_with_violations。",
        "当前仓库保留 5 个场景、25 个视频、完整 JSON/PY/TXT 诊断产物，便于复现实验。"
    ], font_size=14)


def slide_problem(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "核心问题：教学视频不是普通 text-to-video", "代码正确、画面好看、教学目标达成是三件不同的事")
    add_bullets(s, 0.75, 1.65, 5.7, 4.8, [
        "自然语言 concept 往往太抽象，无法直接保证视觉证据。",
        "LLM 容易生成“有对象但不可读、有线条但关系错”的 Blender 场景。",
        "多轮 critic 如果反馈太散，会导致越修越乱：补更多标签、helper，而不是修最小问题。",
        "因此需要从“生成后打补丁”升级为“生成前定义成功状态，并逐层验证”。",
    ], font_size=14)
    add_card(s, 7.0, 1.65, 5.4, 1.25, "失败例子的抽象模式",
             "text object exists ≠ readable text\nray exists ≠ correct incident/normal/reflection relation\nkeyframe exists ≠ continuous visible motion",
             fill=LIGHT, accent=ACCENT)
    add_card(s, 7.0, 3.2, 5.4, 1.25, "框架目标",
             "每一轮都回答：\n缺什么证据？证据来自 critic、AST、contract 还是 preview？下一轮只修哪些高置信目标？",
             fill=PALE_BLUE, accent=BLUE)
    add_card(s, 7.0, 4.75, 5.4, 1.15, "最终态度",
             "能渲染但未达标时，明确输出 best_with_violations，而不是把视频包装成 pass。",
             fill=PALE_GREEN, accent=GREEN)


def slide_pipeline(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "总体架构：从概念到可诊断视频", "当前代码仓库的主链路，不再是早期 W3 版本")
    stages = [
        ("Concept YAML", "普通概念配置\n不要求用户写 metric", BLUE),
        ("Narrative/Profile", "教学节点\n场景语义锚点", TEAL),
        ("Auto Success Spec", "自动软规则\nobject/text/safe-frame", GREEN),
        ("Storyboard/IR", "分镜、对象、相机\nVisual Contract", ORANGE),
        ("scene.py", "LLM Coder\nCompiled Fallback", NAVY),
        ("Checks + Critic", "Verifier/Preview\nClaude+GPT evidence", ACCENT),
        ("Selection", "failure_class\nfinal.mp4", GREEN),
    ]
    x0, y0, w, h, gap = 0.35, 2.0, 1.62, 1.05, 0.20
    for i, (title, body, color) in enumerate(stages):
        x = x0 + i * (w + gap)
        add_card(s, x, y0, w, h, title, body, fill=LIGHT, accent=color)
        if i < len(stages) - 1:
            add_arrow(s, x + w + 0.03, y0 + 0.52, x + w + gap - 0.04, y0 + 0.52)
    add_bullets(s, 0.7, 4.15, 12.0, 1.8, [
        "每个阶段都产生 JSON/PY/TXT artifact：不是黑盒生成，而是可复盘生成。",
        "Deterministic 层负责低成本硬检查；VLM critic 负责视觉观察；cross-reference 负责把视觉问题映射回代码结构。",
        "Compiled fallback 是诊断保底，不允许静默伪装成高质量成功。"
    ], font_size=14)


def slide_feedback(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "反馈闭环：让 critic 的问题真正回流到 coder", "本阶段最重要的架构收口")
    nodes = [
        ("Verifier", 0.9, 1.75, PALE_RED, ACCENT, "syntax/render/security\nmissing render call"),
        ("Contract", 0.9, 3.25, PALE_BLUE, BLUE, "anchors / labels / vectors\nrequired relationships"),
        ("Preview", 0.9, 4.75, LIGHT, NAVY, "motion / visibility / crash\ncheap sampled frames"),
        ("Critic Ensemble", 5.0, 1.75, PALE_GREEN, GREEN, "Claude + GPT\npartial success preserved"),
        ("Cross-ref / Metrics", 5.0, 3.25, LIGHT, TEAL, "critic issue × AST evidence\nAuto Success Spec checks"),
        ("Repair Plan", 5.0, 4.75, PALE_BLUE, BLUE, "≤ 6 high-confidence targets\nminimal repair principle"),
        ("Coder Retry / Best Selection", 9.2, 3.0, PALE_RED, ACCENT, "hard failures before score\nbest_with_violations if needed"),
    ]
    for title, x, y, fill, color, body in nodes:
        add_card(s, x, y, 3.0, 0.95, title, body, fill=fill, accent=color)
    add_arrow(s, 3.95, 2.2, 4.85, 2.2)
    add_arrow(s, 3.95, 3.7, 4.85, 3.7)
    add_arrow(s, 3.95, 5.2, 4.85, 5.2)
    add_arrow(s, 8.1, 2.2, 9.05, 3.25)
    add_arrow(s, 8.1, 3.7, 9.05, 3.55)
    add_arrow(s, 8.1, 5.2, 9.05, 3.85)
    add_text(s, 0.75, 6.25, 12.0, 0.45,
             "关键变化：不再把 critic free text 原样塞回 prompt，而是压缩成少量可执行修复目标。",
             font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_success_spec(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "Auto Success Spec：自动生成软门禁", "用户只写普通 concept YAML，系统自动提取可验证成功证据")
    add_table(
        s, 0.75, 1.65,
        ("Rule kind", "Evidence expected", "Why soft-first"),
        [
            ("object_visible", "anchor exists and is visible", "AST created ≠ critic can see"),
            ("text_readable", "labels/readouts readable", "avoid brittle OCR hard gate"),
            ("stay_in_screen_safe", "important panels not cropped", "common across concepts"),
            ("helper_hidden", "guides do not dominate", "never upgraded to hard"),
            ("animation_coverage", "motion/keyframes cover storyboard", "proxy objects allowed"),
        ],
        (2.8, 4.2, 5.2),
        row_h=0.54,
        font_size=10,
    )
    add_bullets(s, 0.85, 5.05, 11.6, 1.15, [
        "generated rules 初始为 success_soft，不在 iter00 直接 hard fail。",
        "只有连续 critic evidence 且没有 AST 反证，才在当前 run 内升级。",
        "这避免了“规则越多，画面越乱”的 Goodhart 风险。"
    ], font_size=13)


def slide_results(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "当前结果：5 个场景全部可渲染，但不误报 pass", "主指标来自 outputs/<concept>/critic_best.json 与实验文档")
    add_table(
        s, 0.45, 1.55,
        ("Concept", "Best", "Score", "Frame b/w", "Concept b/w", "Final status"),
        RESULT_ROWS,
        (2.35, 0.82, 1.0, 1.25, 1.45, 3.15),
        row_h=0.47,
        font_size=9,
    )
    add_card(s, 0.85, 5.45, 3.6, 1.0, "好消息",
             "5/5 有 final.mp4\nartifact 完整可复盘", fill=PALE_GREEN, accent=GREEN)
    add_card(s, 4.85, 5.45, 3.6, 1.0, "诚实状态",
             "全部 best_with_violations\n不把可渲染包装成达标", fill=PALE_BLUE, accent=BLUE)
    add_card(s, 8.85, 5.45, 3.6, 1.0, "主要短板",
             "复杂语义仍难稳定\nray/normal/label 关系最难", fill=PALE_RED, accent=ACCENT)


def slide_video_grid(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "结果视频展示：当前保留的 5 个 final.mp4", "本页嵌入可播放视频；缩略图作为 poster frame")
    positions = [(0.55, 1.55), (3.15, 1.55), (5.75, 1.55), (8.35, 1.55), (10.95, 1.55)]
    for (concept, zh, note), (x, y) in zip(CONCEPTS, positions):
        thumb = thumbs.get(concept)
        video = ROOT / "outputs" / concept / "final.mp4"
        if video.exists():
            kwargs = {}
            if thumb and thumb.exists():
                kwargs["poster_frame_image"] = str(thumb)
            s.shapes.add_movie(
                str(video),
                Inches(x),
                Inches(y),
                Inches(2.25),
                Inches(1.265),
                mime_type="video/mp4",
                **kwargs,
            )
        elif thumb and thumb.exists():
            s.shapes.add_picture(str(thumb), Inches(x), Inches(y), width=Inches(2.25))
        add_text(s, x, y + 1.45, 2.25, 0.28, zh, font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x, y + 1.74, 2.25, 0.55, note, font_size=8, color=GREY, align=PP_ALIGN.CENTER)
    add_table(
        s, 1.05, 4.7,
        ("File", "Role in presentation"),
        [
            ("final.mp4", "主结果视频；报告表格只统计这一条最终选择"),
            ("final_balanced.mp4", "默认 selector 视角；综合语义和画面"),
            ("final_compliance.mp4 / final_semantic.mp4", "偏规则/语义，用于诊断对比"),
            ("final_aesthetic.mp4", "偏画面质量；不能覆盖 hard failure"),
        ],
        (4.2, 7.2),
        row_h=0.42,
        font_size=10,
    )


def slide_cases(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "案例对比：画面感 vs. 语义压力", "Prism 更适合当前架构；Mirror 暴露 ray/normal 关系短板")
    left = thumbs.get("prism_dispersion_teaching")
    right = thumbs.get("mirror_reflection")
    if left and left.exists():
        s.shapes.add_picture(str(left), Inches(0.75), Inches(1.55), width=Inches(5.2))
    if right and right.exists():
        s.shapes.add_picture(str(right), Inches(7.05), Inches(1.55), width=Inches(5.2))
    add_text(s, 0.75, 4.55, 5.2, 0.35, "prism_dispersion_teaching",
             font_size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_bullets(s, 0.85, 5.0, 5.0, 1.15, [
        "光学台 + RGB rays + prism 是清晰视觉隐喻",
        "score 0.8388，concept block 低",
        "仍需检查内部光路/法线的物理关系",
    ], font_size=11)
    add_text(s, 7.05, 4.55, 5.2, 0.35, "mirror_reflection",
             font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_bullets(s, 7.15, 5.0, 5.0, 1.15, [
        "可渲染但 concept block 高",
        "ray / normal / angle / label 关系最容易错",
        "证明 feedback 能诊断，但 generation 仍是瓶颈",
    ], font_size=11)


def slide_artifacts(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "每次运行留下什么？", "不仅是视频，更是可诊断、可复盘、可继续修的证据链")
    add_card(s, 0.7, 1.6, 3.7, 1.2, "生成产物",
             "final*.mp4\nscene.py / scene.compiled.py\nstoryboard.json / narrative.json",
             fill=PALE_GREEN, accent=GREEN)
    add_card(s, 4.75, 1.6, 3.7, 1.2, "检查产物",
             "scene_verifier*.json\ncontract_validation*.json\npreview_report*.json",
             fill=PALE_BLUE, accent=BLUE)
    add_card(s, 8.8, 1.6, 3.7, 1.2, "反馈产物",
             "critic_iter*.json\nmember_usable_summary\nrepair_plan / cross_reference",
             fill=PALE_RED, accent=ACCENT)
    add_bullets(s, 0.9, 3.65, 11.6, 2.3, [
        "这套 artifact 结构让我们能回答：为什么这一版差？critic 有没有丢问题？retry 是否拿到目标？fallback 是否只是诊断？",
        "最终报告强调当前框架的诚实性：复杂场景仍未严格 pass，但失败原因是可定位的。",
        "提交快照保留当前视频和复盘级 JSON/PY/TXT；旧 W3 图片和过时 concept 已移除。"
    ], font_size=14)


def slide_limitations(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "局限与下一步", "当前框架能诊断问题，但还不能稳定解决所有复杂教学语义")
    add_table(
        s, 0.7, 1.55,
        ("Limitation", "Observed symptom", "Next step"),
        [
            ("Vector/ray semantics", "mirror concept block high", "spatial evidence extractor"),
            ("Temporal continuity", "keyframe samples not enough", "dense motion / frame-diff checks"),
            ("Text readability", "labels may exist but unreadable", "OCR or camera-facing evidence"),
            ("Repair reliability", "LLM may add clutter", "smaller targets + deterministic helpers"),
            ("Generated spec", "soft signals not always enough", "evidence-based escalation only"),
        ],
        (3.0, 4.4, 4.8),
        row_h=0.55,
        font_size=9,
    )
    add_text(s, 0.8, 5.95, 11.7, 0.5,
             "核心结论：下一步不应该无限加 prompt 或手写 metric，而是增强 render evidence extractor。",
             font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_demo_plan(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "10 分钟汇报节奏", "讲清方法，再用视频证明系统当前能做什么、还不能做什么")
    add_table(
        s, 0.9, 1.6,
        ("Time", "Content", "Suggested action"),
        [
            ("0:00-1:00", "选题与问题定义", "说明 Agent + CG 教学动画"),
            ("1:00-3:30", "总体架构与反馈闭环", "讲 pipeline / Success Spec / repair plan"),
            ("3:30-5:30", "当前 5 个结果总览", "展示结果表和视频缩略图"),
            ("5:30-7:30", "播放 2 个视频", "推荐 prism + mirror，对比画面感和语义压力"),
            ("7:30-9:00", "局限与反思", "为什么 best_with_violations 是诚实设计"),
            ("9:00-10:00", "分工与总结", "代码链接、报告、后续方向"),
        ],
        (1.7, 4.4, 6.2),
        row_h=0.58,
        font_size=10,
    )
    add_text(s, 1.0, 6.3, 11.4, 0.35,
             "现场视频路径：outputs/prism_dispersion_teaching/final.mp4 和 outputs/mirror_reflection/final.mp4",
             font_size=12, color=GREEN, align=PP_ALIGN.CENTER)


def slide_close(prs, thumbs):
    s = add_blank(prs)
    add_title(s, "总结", "CG-Tutor 的价值在于“生成 + 验证 + 诚实诊断”的完整链路")
    add_bullets(s, 0.95, 1.65, 11.5, 2.4, [
        "完成了从普通 concept YAML 到 Blender 教学视频的自动化生成框架。",
        "把成功标准显式化：Auto Success Spec、visual contract、failure_class、repair plan。",
        "当前 5 个场景全部可渲染并保留视频，但严格语义仍未全部 pass。",
        "最终结论：复杂教学视频生成不能只靠更长 prompt；需要可验证 evidence 闭环。",
    ], font_size=17)
    add_card(s, 1.2, 5.0, 5.3, 1.05, "代码与报告",
             "https://github.com/AssassinCow/CG-PJ3.git\nreport/main.tex + slides/CG-Tutor.pptx",
             fill=PALE_BLUE, accent=BLUE)
    add_card(s, 7.0, 5.0, 5.3, 1.05, "成员",
             "刘卓鑫 24300240170\n朱劲舟 24300240117",
             fill=PALE_GREEN, accent=GREEN)


def main() -> None:
    thumbs = ensure_thumbnails()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    builders = [
        slide_title,
        slide_requirements,
        slide_problem,
        slide_pipeline,
        slide_feedback,
        slide_success_spec,
        slide_results,
        slide_video_grid,
        slide_cases,
        slide_artifacts,
        slide_limitations,
        slide_demo_plan,
        slide_close,
    ]
    total = len(builders)
    for page, builder in enumerate(builders, start=1):
        builder(prs, thumbs)
        if page > 1:
            add_footer(prs.slides[page - 1], page, total)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT} ({OUT.stat().st_size // 1024} KB, {total} slides)")


if __name__ == "__main__":
    main()
